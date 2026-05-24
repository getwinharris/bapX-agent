#!/usr/bin/env python3
"""bapX PowerPoint tool — Create, read, edit .pptx decks"""
import sys, json, io
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print(json.dumps({"error": "python-pptx not installed. Run: pip install python-pptx"}))
    sys.exit(1)

def cmd_create(args):
    """Create a new presentation with slides from JSON spec"""
    slides_data = json.loads(args.get("slides", "[]"))
    prs = Presentation()
    for s in slides_data:
        layout = prs.slide_layouts[s.get("layout", 6)]  # 6=blank
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        if title and s.get("title"):
            title.text = s["title"]
        for content in s.get("content", []):
            from pptx.util import Inches, Pt
            txBox = slide.shapes.add_textbox(
                Inches(content.get("left", 1)),
                Inches(content.get("top", 1)),
                Inches(content.get("width", 8)),
                Inches(content.get("height", 1))
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = content.get("text", "")
            if content.get("bold"):
                p.runs[0].font.bold = True
            if content.get("size"):
                p.runs[0].font.size = Pt(content["size"])
    path = args.get("output", "presentation.pptx")
    prs.save(path)
    return {"status": "ok", "path": path, "slides": len(slides_data)}

def cmd_read(args):
    """Read a presentation and return slide contents"""
    prs = Presentation(args.get("path", "presentation.pptx"))
    slides = []
    for slide in prs.slides:
        s = {"title": "", "texts": []}
        if slide.shapes.title:
            s["title"] = slide.shapes.title.text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        s["texts"].append(para.text)
        slides.append(s)
    return {"status": "ok", "slides": slides, "count": len(slides)}

def cmd_edit(args):
    """Edit a slide in an existing presentation"""
    path = args.get("path", "presentation.pptx")
    slide_idx = int(args.get("slide", 0))
    new_text = args.get("text", "")
    prs = Presentation(path)
    if slide_idx >= len(prs.slides):
        return {"error": f"Slide index {slide_idx} out of range ({len(prs.slides)} slides)"}
    slide = prs.slides[slide_idx]
    if slide.shapes.title:
        slide.shapes.title.text = new_text
    prs.save(path)
    return {"status": "ok", "path": path, "slide": slide_idx}

if __name__ == "__main__":
    cmd = sys.argv[1] if len(sys.argv) > 1 else "help"
    args = json.loads(sys.argv[2]) if len(sys.argv) > 2 else {}
    cmds = {"create": cmd_create, "read": cmd_read, "edit": cmd_edit}
    if cmd == "help":
        print(json.dumps({"commands": list(cmds.keys()), "usage": "bapX pptx <command> <json-args>"}))
    elif cmd in cmds:
        print(json.dumps(cmds[cmd](args)))
    else:
        print(json.dumps({"error": f"Unknown command: {cmd}"}))
